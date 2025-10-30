import os
import json
import dashscope
from dashscope.api_entities.dashscope_response import Role
from PyPDF2 import PdfReader
import argparse
import sys
from pathlib import Path
import re
from dashscope import MultiModalConversation
import tempfile
from datetime import datetime
import traceback
import logging
# 图像处理库
from PIL import Image

# 导入配置管理器
import config_manager

# 修复：确保所有配置都从小球配置中心读取，不再使用环境变量

# 创建运行日志目录 - 按日期创建日志文件
today = datetime.now().strftime('%Y%m%d')
log_dir = "run-log"
os.makedirs(log_dir, exist_ok=True)

# 配置日志 - 按天记录到同一天的日志文件
log_filename = f"{log_dir}/{today}.log"

# 为Windows系统设置控制台编码
try:
    if sys.platform.startswith('win'):
        # 设置控制台编码为UTF-8
        import ctypes
        kernel32 = ctypes.windll.kernel32
        kernel32.SetConsoleCP(65001)  # 设置控制台输入编码为UTF-8
        kernel32.SetConsoleOutputCP(65001)  # 设置控制台输出编码为UTF-8
        
        # 同时设置标准输出和标准错误的编码
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
except Exception as e:
    print(f"设置控制台编码时出错: {str(e)}")

# 自定义StreamHandler类来处理编码问题
class UnicodeStreamHandler(logging.StreamHandler):
    def __init__(self, stream=None):
        super().__init__(stream)
        self.encoding = 'utf-8'
    
    def emit(self, record):
        try:
            # 确保消息是UTF-8编码的字符串
            msg = self.format(record)
            if isinstance(msg, str):
                msg = msg.encode('utf-8', errors='replace').decode('utf-8')
            
            stream = self.stream
            stream.write(msg)
            stream.write(self.terminator)
            self.flush()
        except RecursionError:
            raise
        except Exception:
            self.handleError(record)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_filename, encoding='utf-8'),
        UnicodeStreamHandler()  # 使用自定义的StreamHandler
    ],
    force=True  # 强制重新配置日志系统
)

logger = logging.getLogger(__name__)

def log_info(message):
    """记录信息日志"""
    logger.info(message)

def log_error(message):
    """记录错误日志"""
    logger.error(message)

def log_debug(message):
    """记录调试日志"""
    logger.debug(message)

def read_pdf(file_path):
    """
    读取PDF文件内容
    """
    try:
        log_info(f"开始读取PDF文件: {file_path}")
        reader = PdfReader(file_path)
        text_content = ""
        for page in reader.pages:
            text_content += page.extract_text() + "\n"
        log_info(f"PDF文件读取完成，总字符数: {len(text_content)}")
        return text_content
    except Exception as e:
        error_msg = f"读取PDF文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None

def read_markdown(file_path):
    """
    读取markdown文件内容
    """
    try:
        log_info(f"开始读取markdown文件: {file_path}")
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        log_info(f"markdown文件读取完成，总字符数: {len(content)}")
        return content
    except Exception as e:
        error_msg = f"读取markdown文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None

def read_python_file(file_path):
    """
    读取Python代码文件内容，保留代码结构和注释
    """
    try:
        log_info(f"开始读取Python文件: {file_path}")
        
        # 使用AST模块解析Python代码获取结构信息
        import ast
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
        # 尝试解析代码以验证语法正确性
        try:
            tree = ast.parse(content)
            word_length = len(content.split('\n'))
            # 提取代码结构信息
            functions = []
            classes = []
            
            # 遍历AST提取函数和类
            for node in ast.walk(tree):
                if isinstance(node, ast.FunctionDef):
                    functions.append(node.name)
                elif isinstance(node, ast.ClassDef):
                    classes.append(node.name)

            # 生成代码结构摘要
            structure_info = f"""
# Python文件结构摘要
## 函数列表: {', '.join(functions) if functions else '无'}
## 类列表: {', '.join(classes) if classes else '无'}
## 总行数: {word_length}
## 字符数: {len(content)}

# 原始代码内容:
"""
            
            # 组合结构摘要和原始代码
            full_content = structure_info + content
            
        except SyntaxError:
            # 如果语法解析失败，仍然返回原始内容但添加警告
            log_info(f"Python文件 {file_path} 包含语法错误，但仍将处理原始内容")
            full_content = f"""
# 警告: 此Python文件包含语法错误
# 文件名: {os.path.basename(file_path)}

# 原始代码内容:
""" + content
        
        log_info(f"Python文件读取完成，总字符数: {len(full_content)}")
        return full_content
    except Exception as e:
        error_msg = f"读取Python文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None

def compress_image(image_path, max_size_mb=5, quality=95):
    """
    压缩图像文件以减小大小但保持质量，特别优化文字清晰度
    
    压缩机制说明：
    1. 先尝试仅通过调整JPEG/PNG压缩质量进行无损或近无损压缩
    2. 仅当质量压缩后仍超限时，才考虑降低分辨率
    3. 使用LANCZOS重采样算法保证文字边缘清晰
    4. 对文字内容优先保持水平和垂直分辨率
    
    Args:
        image_path: 原始图像路径
        max_size_mb: 最大文件大小（MB）
        quality: 压缩质量 (0-100)，默认95保持较高质量
    
    Returns:
        str: 压缩后的临时图像路径，如果不需要压缩则返回原路径
    """
    import os
    from PIL import Image, ImageFilter
    import tempfile
    
    # 检查文件大小
    file_size_mb = os.path.getsize(image_path) / (1024 * 1024)
    log_info(f"原始图像大小: {file_size_mb:.2f}MB")
    
    # 如果文件大小在限制内，直接返回原路径
    if file_size_mb <= max_size_mb:
        log_info(f"图像大小在限制范围内，无需压缩")
        return image_path
    
    log_info(f"图像大小超过限制，开始优化压缩...")
    
    try:
        # 创建临时文件
        file_ext = os.path.splitext(image_path)[1]
        temp_fd, temp_path = tempfile.mkstemp(suffix=file_ext)
        os.close(temp_fd)
        
        # 打开图像
        with Image.open(image_path) as img:
            width, height = img.size
            original_format = img.format
            
            # 第一步：先尝试无损或近无损压缩，不改变分辨率
            img.save(temp_path, format=original_format, quality=quality, optimize=True)
            compressed_size_mb = os.path.getsize(temp_path) / (1024 * 1024)
            
            if compressed_size_mb <= max_size_mb:
                log_info(f"仅通过调整压缩质量完成优化: {file_size_mb:.2f}MB → {compressed_size_mb:.2f}MB，保持原始分辨率 {width}x{height}")
                return temp_path
            
            # 第二步：如果质量压缩后仍超限，尝试降低分辨率
            # 对于可能包含文字的图像，优先考虑更高的质量设置
            target_size_mb = max_size_mb * 0.9  # 留出余量
            
            # 计算目标分辨率，基于文件大小比例的平方根，这样可以保持宽高比
            size_ratio = (target_size_mb / compressed_size_mb) ** 0.5
            # 对于文字内容，不要过度降低分辨率
            min_ratio = 0.5  # 至少保持50%的原始分辨率
            ratio = max(size_ratio, min_ratio)
            
            new_width = int(width * ratio)
            new_height = int(height * ratio)
            
            # 确保分辨率不会过高或过低
            max_dimension = 3000  # 最大允许分辨率
            min_dimension = 800    # 最小允许分辨率（保证文字可读性）
            
            new_width = max(min(new_width, max_dimension), min_dimension)
            new_height = max(min(new_height, max_dimension), min_dimension)
            
            log_info(f"需要调整分辨率: 从 {width}x{height} 到 {new_width}x{new_height}，保持质量设置为 {quality}")
            
            # 使用LANCZOS算法（最高质量）进行重采样，特别适合保持文字清晰度
            resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # 对于JPEG格式，添加少量锐化以保持文字边缘清晰
            if original_format == 'JPEG':
                resized_img = resized_img.filter(ImageFilter.SHARPEN)
            
            # 保存压缩后的图像
            # 对于可能包含文字的图像，优先考虑PNG格式以获得更好的文字清晰度
            if original_format in ['JPEG', 'JPG'] and compressed_size_mb > 2:  # 如果JPEG压缩不够好，尝试PNG
                try:
                    resized_img.save(temp_path, format='PNG', optimize=True)
                    png_size_mb = os.path.getsize(temp_path) / (1024 * 1024)
                    if png_size_mb <= max_size_mb:
                        log_info(f"已转换为PNG格式以优化文字清晰度: {file_size_mb:.2f}MB → {png_size_mb:.2f}MB")
                        return temp_path
                except Exception as png_error:
                    log_info(f"PNG转换失败，继续使用原格式: {str(png_error)}")
            
            # 使用原格式保存
            resized_img.save(temp_path, format=original_format, quality=quality, optimize=True)
            
            # 检查最终压缩后的大小
            final_size_mb = os.path.getsize(temp_path) / (1024 * 1024)
            log_info(f"图像压缩完成: {file_size_mb:.2f}MB → {final_size_mb:.2f}MB，分辨率: {new_width}x{new_height}")
            
            return temp_path
    except Exception as e:
        log_error(f"压缩图像时出错: {str(e)}")
        # 如果压缩失败，尝试使用原图
        return image_path

def read_image_file(file_path, api_key):
    """
    读取图像文件并使用DashScope API进行识别
    """
    import os
    log_info(f"开始处理图像文件: {file_path}")
    temp_compressed_path = None
    
    try:
        # 先压缩图像以避免文件过大
        temp_compressed_path = compress_image(file_path)
        
        # 使用recognize_image_with_dashscope函数识别图像
        description = recognize_image_with_dashscope(
            api_key,
            temp_compressed_path,
            "请详细描述这张图片的内容，包括其中的关键信息、文字、数据、图表、表格或其他重要元素。如果图片包含文字，请尽可能准确地提取，图像可能不清晰，帮我推断一下主题和提炼图片所包含的主要信息。"
        )
        
        if description:
            # 构建markdown格式的图像识别结果
            file_name = os.path.basename(file_path)
            image_content = f"# 图像识别结果: {file_name}\n\n"
            image_content += f"## 图像基本信息\n\n"
            image_content += f"- **文件名**: {file_name}\n"
            
            # 获取图像尺寸信息
            try:
                from PIL import Image
                with Image.open(file_path) as img:
                    width, height = img.size
                    image_content += f"- **尺寸**: {width} x {height} 像素\n"
                    image_content += f"- **格式**: {img.format}\n"
                    
                    # 获取文件大小
                    file_size = os.path.getsize(file_path)
                    image_content += f"- **文件大小**: {file_size/1024:.2f} KB\n"
            except Exception as img_info_error:
                log_error(f"获取图像信息时出错: {str(img_info_error)}")
                image_content += "- **尺寸**: 无法获取\n"
                image_content += "- **格式**: 无法获取\n"
                
            image_content += "\n## 图像内容描述\n\n"
            image_content += description
            
            # 添加图像引用（在markdown中显示图像）
            image_content += f"\n\n## 原始图像\n\n!['{file_name}']({file_path})"
            
            log_info(f"图像文件处理成功，识别完成")
            return image_content
        else:
            log_error(f"图像识别失败: {file_path}")
            return None
    except Exception as e:
        error_msg = f"处理图像文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None
    finally:
        # 清理临时压缩文件
        if temp_compressed_path and temp_compressed_path != file_path:
            try:
                if os.path.exists(temp_compressed_path):
                    os.remove(temp_compressed_path)
                    log_info(f"已清理临时压缩文件: {temp_compressed_path}")
            except Exception as cleanup_error:
                log_error(f"清理临时文件时出错: {str(cleanup_error)}")

def read_ppt(file_path):
    """
    读取PPT文件内容，提取文本和图片
    """
    try:
        log_info(f"开始读取PPT文件: {file_path}")
        
        # 先检查pptx模块是否已安装
        try:
            import pptx
            log_info(f"python-pptx模块已安装，版本: {pptx.__version__}")
        except ImportError:
            log_error("python-pptx模块未安装，尝试自动安装...")
            try:
                # 尝试自动安装依赖
                import subprocess
                subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx>=0.6.23'])
                log_info("python-pptx模块安装成功")
            except Exception as install_error:
                log_error(f"自动安装python-pptx失败: {str(install_error)}")
                log_error("请手动运行: pip install python-pptx>=0.6.23")
                raise ImportError("缺少python-pptx依赖，请安装: pip install python-pptx>=0.6.23")
        
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import os
        
        prs = Presentation(file_path)
        content = {"slides": [], "images": []}
        
        for i, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": i + 1,
                "text": [],
                "images": []
            }
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content["text"].append(shape.text)
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        text = ""
                        for paragraph in shape.text_frame.paragraphs:
                            text += "".join(run.text for run in paragraph.runs)
                        if text.strip():
                            slide_content["text"].append(text)
            
            # 提取图片
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    # 保存图片到临时文件
                    image_filename = f"slide_{i+1}_image_{len(slide_content['images'])+1}.png"
                    image_path = os.path.join(tempfile.gettempdir(), image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image.blob)
                    slide_content["images"].append(image_path)
                    content["images"].append(image_path)
            
            content["slides"].append(slide_content)
        
        log_info(f"PPT文件读取完成，共 {len(content['slides'])} 张幻灯片")
        return content
    except ImportError as e:
        # 专门处理模块未找到的情况
        error_msg = f"读取PPT文件时缺少依赖: {str(e)}"
        log_error(error_msg)
        log_error("请运行: pip install python-pptx")
        raise ImportError(error_msg) from e  # 重新抛出异常以确保被外层捕获
    except Exception as e:
        error_msg = f"读取PPT文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        raise Exception(error_msg) from e  # 重新抛出异常以确保被外层捕获

def format_ppt_content_for_markdown(ppt_content):
    """
    将PPT内容格式化为markdown格式
    """
    log_info("开始格式化PPT内容为markdown格式")
    md_content = f"# PPT内容整理\n\n"
    md_content += f"总共 {len(ppt_content['slides'])} 张幻灯片\n\n"
    
    for slide in ppt_content["slides"]:
        md_content += f"## 第 {slide['slide_number']} 张幻灯片\n\n"
        
        for text in slide["text"]:
            if text.strip():
                md_content += f"{text}\n\n"
        
        for img_path in slide["images"]:
            img_name = os.path.basename(img_path)
            md_content += f"![{img_name}]({img_path})\n\n"
    
    log_info(f"PPT内容格式化完成，总字符数: {len(md_content)}")
    return md_content

def extract_images_from_markdown(markdown_content, base_path):
    """
    从markdown内容中提取图片路径
    """
    log_info("开始从markdown内容中提取图片路径")
    # 匹配markdown中的图片语法 ![alt text](image_path)
    image_pattern = r'!\[.*?\]\((.*?)\)'
    matches = re.findall(image_pattern, markdown_content)
    
    image_paths = []
    for match in matches:
        # 处理相对路径和绝对路径
        if match.startswith('http://') or match.startswith('https://'):
            # 远程图片URL
            image_paths.append(match)
        else:
            # 本地图片路径，转换为绝对路径
            img_path = Path(base_path).parent / match
            if img_path.exists():
                image_paths.append(str(img_path.resolve()))
            else:
                warning_msg = f"警告: 图片文件不存在 - {img_path}"
                log_info(warning_msg)
    
    log_info(f"提取到 {len(image_paths)} 个图片路径")
    return image_paths

def recognize_image_with_dashscope(api_key, image_path, custom_prompt="请详细描述这张图片的内容"):
    """
    使用DashScope视觉模型识别图片
    """
    log_info(f"开始识别图片: {image_path}")
    log_info(f"使用提示词: {custom_prompt}")
    dashscope.api_key = api_key
    
    # 加载配置以获取图像模型名称
    config = config_manager.load_config()
    image_model = config.get('image_model', 'qwen-vl-plus')
    
    try:
        # 构建消息，包含图片和描述请求
        if image_path.startswith('http://') or image_path.startswith('https://'):
            # 远程图片URL
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"image": image_path},
                        {"text": custom_prompt}
                    ]
                }
            ]
        else:
            # 本地图片文件
            messages = [
                {
                    "role": "user",
                    "content": [
                        {"image": f"file://{image_path}"},
                        {"text": custom_prompt}
                    ]
                }
            ]
        
        log_info(f"调用DashScope视觉模型API...")
        log_info(f"使用图像模型: {image_model}")
        
        # 使用从配置管理器获取的图像模型
        response = MultiModalConversation.call(
            model=image_model,
            messages=messages
        )
        
        log_info(f"API响应状态码: {response.status_code}")
        
        if response.status_code == 200:
            # 提取识别结果
            result = response.output.choices[0].message.content
            log_info(f"图片识别API调用成功")
            log_debug(f"API响应内容: {result}")
            
            # 提取并统计token用量
            image_token_usage = 0
            if hasattr(response, 'usage') and response.usage:
                input_tokens = response.usage.get('input_tokens', 0)
                output_tokens = response.usage.get('output_tokens', 0)
                total_tokens = input_tokens + output_tokens
                image_token_usage = total_tokens
                log_info(f"图片识别Token用量统计: 输入{input_tokens} + 输出{output_tokens} = 总计{total_tokens}")
            else:
                log_info("未在图像识别API响应中找到token用量信息")
            
            # 输出图片识别的token用量，以便web_app.py能够捕获和记录
            print(f"IMAGE_TOKEN_USAGE:{image_token_usage}")
            
            # 处理返回的多模态内容
            if isinstance(result, list):
                # 尝试从列表中提取文本内容
                for item in result:
                    if isinstance(item, dict):
                        # 检查是否有text字段
                        if 'text' in item:
                            log_info(f"图片识别结果(文本): {item.get('text', '')[:100]}...")  # 只记录前100个字符
                            return item.get('text', '')
                        # 尝试直接获取文本（不通过type检查）
                        elif 'content' in item:
                            log_info(f"图片识别结果(内容): {item.get('content', '')[:100]}...")  # 只记录前100个字符
                            return item.get('content', '')
                # 如果列表中没有找到有效文本，尝试将整个列表转为字符串
                result_str = str(result)
                log_info(f"图片识别结果(列表转换): {result_str[:100]}...")
                return result_str
            elif isinstance(result, dict):
                # 如果是字典，直接提取其中的文本内容
                if 'text' in result:
                    log_info(f"图片识别结果(字典文本): {result.get('text', '')[:100]}...")
                    return result.get('text', '')
                # 尝试其他可能的字段名
                elif 'content' in result:
                    log_info(f"图片识别结果(字典内容): {result.get('content', '')[:100]}...")
                    return result.get('content', '')
                else:
                    # 将字典转为字符串
                    result_str = str(result)
                    log_info(f"图片识别结果(字典转换): {result_str[:100]}...")
                    return result_str
            elif isinstance(result, str):
                # 如果直接是字符串，直接返回
                log_info(f"图片识别结果(直接字符串): {result[:100]}...")
                return result
            else:
                # 其他类型转为字符串
                result_str = str(result)
                log_info(f"图片识别结果(其他类型转换): {result_str[:100]}...")
                return result_str
        else:
            error_msg = f"图像识别API调用失败: {response.code} - {response.message}"
            log_error(error_msg)
            return None
    except Exception as e:
        error_msg = f"图像识别时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None

def process_markdown_with_images(api_key, markdown_content, base_path, user_prompt, config):
    """
    处理包含图片的markdown文件，对图片进行识别并整合内容
    """
    log_info("开始处理markdown文件中的图片...")
    log_info(f"用户提示词: {user_prompt}")
    
    # 提取图片路径
    image_paths = extract_images_from_markdown(markdown_content, base_path)
    
    # 存储图片识别结果
    image_descriptions = {}
    
    # 对每张图片进行识别
    for img_path in image_paths:
        log_info(f"正在识别图片: {img_path}")
        description = recognize_image_with_dashscope(
            api_key, 
            img_path, 
            "请详细描述这张图片的内容，包括其中的关键信息、文字、数据或其他重要元素。"
        )
        if description:
            image_descriptions[img_path] = description
            log_info(f"图片识别完成: {img_path}")
        else:
            log_error(f"图片识别失败: {img_path}")
    
    # 将图片描述整合到原始markdown内容中
    enhanced_content = markdown_content
    
    for img_path, description in image_descriptions.items():
        # 在图片下方添加描述
        image_ref = f"!["
        # 找到图片引用的位置并添加描述
        if img_path.startswith('http'):
            # 处理URL图片
            image_tag = f"!["
            pattern = r'!\[.*?\]\(' + re.escape(img_path) + r'\)'
        else:
            # 处理本地图片
            original_img_path = Path(img_path).name
            pattern = r'!\[.*?\][^)]*' + re.escape(original_img_path) + r'[^)]*\)'
        
        # 替换markdown，添加图片描述
        enhanced_content = re.sub(
            pattern, 
            f'![图示]({img_path})\n\n**图片描述**: {description}\n\n', 
            enhanced_content
        )
    
    # 如果没有图片，直接返回原始内容
    if not image_paths:
        # 仍然返回原始内容，但告知用户没有找到图片
        log_info("未在markdown中找到图片")
        return markdown_content
    else:
        log_info("markdown图片处理完成")
        return enhanced_content

def call_dashscope_api(api_key, content, user_prompt, config, file_type="pdf"):
    """
    调用DashScope API处理内容
    """
    log_info(f"开始调用DashScope API处理{file_type}内容...")
    log_info(f"内容长度: {len(content)} 字符")
    log_info(f"用户提示词: {user_prompt}")
    
    # 修复：使用传入的api_key参数，而不是从环境变量获取
    dashscope.api_key = api_key
    
    # 修复：使用传入的config参数，确保与config_manager保持一致
    
    # 使用统一的通用提示词处理所有文件类型
    default_prompt = config.get("default_prompt", "")
    
    # 如果是Python文件，可以添加额外的代码分析提示
    if file_type == "python":
        python_specific = "\n\n特别注意：如果内容包含代码，请分析代码结构和功能，提取关键组件，总结核心逻辑。"
    else:
        python_specific = ""
    
    # 如果没有配置提示词，使用空字符串
    default_prompt = default_prompt or ""
    
    full_prompt = f"{default_prompt}{python_specific}\n\n额外要求: {user_prompt}\n\n内容如下:\n\n{content}"
    
    log_info(f"构建的完整提示词长度: {len(full_prompt)} 字符")
    log_debug(f"完整提示词内容: {full_prompt[:500]}...")  # 只记录前500个字符
    
    messages = [
        {"role": "system", "content": "你是一个专业的文档处理助手，擅长将各种内容整理成结构化markdown格式文档"},
        {"role": "user", "content": full_prompt}
    ]
    
    # 初始化token用量
    total_tokens = 0
    
    try:
        # 从配置中获取模型名称，如果未配置则使用默认值
        text_model = config.get('text_model', 'qwen-plus')
        
        response = dashscope.Generation.call(
            model=text_model,
            messages=messages,
            result_format='message'
        )
        
        log_info(f"API响应状态码: {response.status_code}")
        
        if response.status_code == 200:
            result = response.output.choices[0].message.content
            log_info(f"API调用成功，返回内容长度: {len(result) if result else 0} 字符")
            log_debug(f"API返回内容: {result[:500] if result else 'None'}...")  # 只记录前500个字符
            
            # 提取token用量信息
            if hasattr(response, 'usage') and response.usage:
                input_tokens = response.usage.get('input_tokens', 0)
                output_tokens = response.usage.get('output_tokens', 0)
                total_tokens = input_tokens + output_tokens
                log_info(f"Token用量统计: 输入{input_tokens} + 输出{output_tokens} = 总计{total_tokens}")
            else:
                log_info("未在API响应中找到token用量信息")
                
            # 输出token用量，以便web_app.py捕获
            print(f"TOKEN_USAGE:{total_tokens}")
            
            return result
        else:
            error_msg = f"API调用失败: {response.code} - {response.message}"
            log_error(error_msg)
            return None
    except Exception as e:
        error_msg = f"API调用时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return None

def process_ppt_with_images(api_key, content, base_path, user_prompt, config):
    """
    处理PPT中的图片，对图片进行识别并整合内容
    """
    log_info("开始处理PPT内容中的图片...")
    log_info(f"用户提示词: {user_prompt}")
    
    # 提取图片路径（从markdown格式的图片语法中提取）
    image_paths = extract_images_from_markdown(content, base_path)
    
    # 存储图片识别结果
    image_descriptions = {}
    
    # 对每张图片进行识别
    for img_path in image_paths:
        log_info(f"正在识别图片: {img_path}")
        description = recognize_image_with_dashscope(
            api_key, 
            img_path, 
            "请详细描述这张图片的内容，包括其中的关键信息、文字、数据或其他重要元素。"
        )
        if description:
            image_descriptions[img_path] = description
            log_info(f"图片识别完成: {img_path}")
        else:
            log_error(f"图片识别失败: {img_path}")
    
    # 将图片描述整合到原始内容中
    enhanced_content = content
    
    for img_path, description in image_descriptions.items():
        # 使用正则表达式找到图片引用并添加描述
        original_img_path = Path(img_path).name
        pattern = r'!\[.*?\][^)]*' + re.escape(original_img_path) + r'[^)]*\)'
        
        # 替换markdown，添加图片描述
        enhanced_content = re.sub(
            pattern, 
            f'![图示]({img_path})\n\n**图片描述**: {description}\n\n', 
            enhanced_content
        )
    
    # 如果没有图片，直接返回原始内容
    if not image_paths:
        log_info("未在PPT内容中找到图片")
        return content
    else:
        log_info("PPT图片处理完成")
        return enhanced_content

def save_markdown(content, output_path):
    """
    保存内容到markdown文件
    """
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
        log_info(f"文档已保存到: {output_path}")
        return True
    except Exception as e:
        error_msg = f"保存文件时出错: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        return False

# 使用config_manager模块加载配置，不再使用自定义load_config函数

def main():
    """
    主函数：处理命令行参数并调用相应的功能函数
    """
    try:
        log_info("开始执行PDF转知识库程序")
        
        parser = argparse.ArgumentParser(description='将PDF、PPT、markdown或Python文件通过大模型API转换为格式化的知识库文档')
        parser.add_argument('input_files', nargs='*', help='输入文件路径列表(PDF、PPT、markdown或Python)')  # 修改为支持多个文件，使用*使其可选
        parser.add_argument('--prompt', '-p', default='', help='额外的个性化提示词')
        parser.add_argument('--output', '-o', help='输出文件路径')
        parser.add_argument('--api-key', help='DashScope API Key')
        
        args = parser.parse_args()
        
        log_info(f"输入参数: input_files={args.input_files}, prompt={args.prompt}, output={args.output}")
        
        # 检查是否提供了输入文件
        if not args.input_files:
            error_msg = "错误: 请提供至少一个输入文件"
            log_error(error_msg)
            parser.print_help()
            sys.exit(1)
        
        # 检查所有输入文件是否存在
        for input_path in args.input_files:
            if not os.path.exists(input_path):
                error_msg = f"错误: 输入文件不存在 - {input_path}"
                log_error(error_msg)
                sys.exit(1)
        
        # 加载配置
        config = config_manager.load_config()
        
        # 获取API KEY，优先级：命令行参数 > 配置管理器
        api_key = args.api_key or config.get('api_key', '')
        
        if not api_key:
            error_msg = "错误: 未提供API KEY，请在系统配置中设置API密钥"
            log_error(error_msg)
            sys.exit(1)
        
        log_info("API KEY已验证")
        
        # 处理多个文件内容
        all_contents = []
        total_files = len(args.input_files)
        has_python_file = False
        has_image_file = False
        
        for i, input_path in enumerate(args.input_files):
            # 判断文件类型
            file_ext = Path(input_path).suffix.lower()
            if file_ext not in ['.pdf', '.md', '.markdown', '.ppt', '.pptx', '.py', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                error_msg = f"错误: 不支持的文件格式 - {file_ext}. 支持的格式: PDF, MD, MARKDOWN, PPT, PPTX, PY, JPG, JPEG, PNG, GIF, BMP, TIFF, WEBP"
                log_error(error_msg)
                sys.exit(1)
            
            log_info(f"处理第 {i+1}/{total_files} 个文件: {input_path} (类型: {file_ext})")
            
            # 根据文件类型处理内容
            if file_ext in ['.pdf']:
                # PDF处理
                log_info("开始处理PDF文件...")
                content = read_pdf(input_path)
                file_type = "pdf"
            elif file_ext in ['.md', '.markdown']:
                # Markdown处理
                log_info("开始处理markdown文件...")
                content = read_markdown(input_path)
                if content:
                    # 处理markdown中的图片
                    log_info("处理markdown中的图片...")
                    content = process_markdown_with_images(api_key, content, input_path, args.prompt, config)
                file_type = "markdown"
            elif file_ext in ['.ppt', '.pptx']:
                # PPT处理
                log_info("开始处理PPT文件...")
                content = read_ppt(input_path)
                if content:
                    # 格式化PPT内容为markdown
                    content = format_ppt_content_for_markdown(content)
                    # 处理PPT中的图片
                    log_info("处理PPT中的图片...")
                    content = process_ppt_with_images(api_key, content, input_path, args.prompt, config)
                file_type = "ppt"
            elif file_ext in ['.py']:
                # Python代码处理
                log_info("开始处理Python代码文件...")
                content = read_python_file(input_path)
                file_type = "python"
                has_python_file = True
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                # 图像文件处理
                log_info("开始处理图像文件...")
                content = read_image_file(input_path, api_key)
                file_type = "image"
                has_image_file = True
            
            if not content:
                # # 对于图像文件，如果识别失败，直接结束任务
                # if file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                #     error_msg = f"错误: 图像识别失败，任务直接结束 - {input_path}"
                #     log_error(error_msg)
                #     print(f"IMAGE_RECOGNITION_FAILED")  # 输出特殊标记供web_app.py捕获
                #     sys.exit(1)
                # else:
                #     # 对于非图像文件，仍然跳过当前文件继续处理
                #     error_msg = f"错误: 无法读取文件内容 - {input_path}"
                #     log_error(error_msg)
                #     continue  # 跳过这个文件，继续处理其他文件
                error_msg = f"错误: 图像识别失败，任务直接结束 - {input_path}"
                log_error(error_msg)
                print(f"IMAGE_RECOGNITION_FAILED")  # 输出特殊标记供web_app.py捕获
                sys.exit(1)
            # 为每个文件内容添加文件名标识
            file_content_header = f"\n# 文件: {os.path.basename(input_path)}\n\n"
            all_contents.append(file_content_header + content)
            
            log_info(f"文件 {input_path} 内容读取成功，字符数: {len(content)}")
        
        # 检查是否有任何内容被成功读取
        if not all_contents:
            error_msg = "错误: 没有任何文件内容被成功读取"
            log_error(error_msg)
            sys.exit(1)
        
        # 合并所有文件内容
        combined_content = "\n".join(all_contents)
        
        log_info(f"所有文件内容合并成功，总字符数: {len(combined_content)}")
        
        # 如果内容过长，进行分段处理提示
        if len(combined_content) > 30000:  # 模型输入限制调整
            log_info("警告: 内容较长，可能超出API限制，正在发送请求...")
        
        # 为Python文件和图像文件添加特殊提示词
        final_prompt = args.prompt
        
        # 如果有Python文件且用户未提供提示词，添加Python相关提示
        if has_python_file and not args.prompt:
            python_prompt_suffix = "\n\n特别注意：对于Python代码，请分析代码结构和功能，提取关键组件，总结核心逻辑，解释主要函数和类的作用。"
            final_prompt = python_prompt_suffix
        
        # 如果有图像文件，添加图像处理相关提示
        if has_image_file:
            image_prompt_suffix = "\n\n对于图像识别结果，请确保详细总结图像中的所有关键信息，包括文本内容、图表数据、视觉元素等。如果有多张图像，请分别详细分析每张图像的内容。"
            if final_prompt:
                final_prompt += image_prompt_suffix
            else:
                final_prompt = image_prompt_suffix
        
        # 调用API处理内容
        log_info("正在调用大模型API处理合并内容...")
        result = call_dashscope_api(api_key, combined_content, final_prompt, config, "multiple_files")
        
        if not result:
            error_msg = "错误: API调用失败，无法生成markdown文档"
            log_error(error_msg)
            sys.exit(1)
        
        # 确定输出路径
        if args.output:
            output_path = args.output
        else:
            # 默认输出路径为第一个文件名+timestamp+processed.md后缀，存放在配置的输出目录
            first_file_name = os.path.splitext(os.path.basename(args.input_files[0]))[0]
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_dir = config.get("output_dir", "./output")
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{first_file_name}_multiple_{timestamp}_processed.md")
        
        log_info(f"输出文件路径: {output_path}")
        
        # 保存结果
        if save_markdown(result, output_path):
            log_info("批量处理完成！")
            print(f"转换完成! 结果已保存到: {output_path}")
        else:
            error_msg = "保存文件失败"
            log_error(error_msg)
            sys.exit(1)
    except Exception as e:
        error_msg = f"处理过程中发生未捕获的异常: {str(e)}"
        log_error(error_msg)
        log_error(f"详细错误信息: {traceback.format_exc()}")
        print(f"错误: 处理失败: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()